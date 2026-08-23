# -*- coding: utf-8 -*-
"""
SlideCast 퀵 가이드 덱 생성기

발표자료 작성법·생성 프로세스·아키텍처·실행 방법을 담은 14장짜리 PPTX를 만든다.
각 슬라이드에는 나레이션용 노트가 들어 있어, 생성된 파일을 그대로 SlideCast
파이프라인에 넣으면 이 가이드를 설명하는 영상이 나온다.

    python scripts/make_guide_deck.py docs/guide/SlideCast_퀵가이드.pptx
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG      = RGBColor(0x11, 0x14, 0x1E)
CARD    = RGBColor(0x1B, 0x20, 0x2E)
CARD2   = RGBColor(0x22, 0x2A, 0x3C)
WHITE   = RGBColor(0xF2, 0xF5, 0xFA)
MUTED   = RGBColor(0x9A, 0xA6, 0xBD)
ACCENT  = RGBColor(0x7C, 0xA9, 0xF0)
GREEN   = RGBColor(0x66, 0xD9, 0xA0)
AMBER   = RGBColor(0xF2, 0xC0, 0x5C)
CORAL   = RGBColor(0xF2, 0x8B, 0x82)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5

def new(notes):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG
    s.notes_slide.notes_text_frame.text = notes
    return s

def txt(s, x, y, w, h, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        space=0, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.alignment = align
        p.space_after = Pt(space)
        p.font.size = Pt(size); p.font.bold = bold; p.font.italic = italic
        p.font.color.rgb = color
    return tb

def box(s, x, y, w, h, fill=CARD, line=None, radius=0.04):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    else:    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh

def head(s, kicker, title, no):
    txt(s, 0.9, 0.62, 11.5, 0.32, kicker, 13, ACCENT, bold=True)
    txt(s, 0.9, 1.02, 11.5, 0.85, title, 34, WHITE, bold=True)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.96), Inches(1.1), Inches(0.045))
    r.fill.solid(); r.fill.fore_color.rgb = ACCENT; r.line.fill.background(); r.shadow.inherit = False
    txt(s, W-1.9, 6.72, 1.0, 0.3, f"{no:02d}", 12, MUTED, align=PP_ALIGN.RIGHT)
    txt(s, 0.9, 6.72, 6.0, 0.3, "SlideCast 퀵 가이드", 12, RGBColor(0x4C,0x56,0x6B))

def card(s, x, y, w, h, title, body, accent=ACCENT, tsize=17, bsize=13, fill=CARD):
    box(s, x, y, w, h, fill=fill)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y+0.22), Inches(0.05), Inches(h-0.44))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit=False
    txt(s, x+0.36, y+0.3, w-0.66, 0.4, title, tsize, WHITE, bold=True)
    txt(s, x+0.36, y+0.3+0.44, w-0.66, h-0.9, body, bsize, MUTED, space=5)

def chip(s, x, y, w, h, label, sub="", fill=CARD2, color=WHITE, edge=None):
    box(s, x, y, w, h, fill=fill, line=edge)
    if sub:
        txt(s, x+0.14, y+0.22, w-0.28, 0.34, label, 13.5, color, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x+0.14, y+0.62, w-0.28, 0.34, sub, 10.5, MUTED, align=PP_ALIGN.CENTER)
    else:
        txt(s, x+0.14, y+(h-0.28)/2, w-0.28, 0.34, label, 13.5, color, bold=True, align=PP_ALIGN.CENTER)

def arrow(s, x, y, w=0.34):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.2))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0x44,0x50,0x68); a.line.fill.background(); a.shadow.inherit=False

def code(s, x, y, w, lines, size=12.5):
    h = 0.30*len(lines) + 0.44
    box(s, x, y, w, h, fill=RGBColor(0x0B,0x0E,0x16), line=RGBColor(0x2A,0x33,0x47))
    tb = s.shapes.add_textbox(Inches(x+0.26), Inches(y+0.2), Inches(w-0.5), Inches(h-0.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.space_after = Pt(3)
        p.font.size = Pt(size); p.font.name = "Menlo"
        p.font.color.rgb = GREEN if ln.startswith("$") else MUTED
    return h

# ───────────────────────── 01 표지
s = new("""키메시지: 표지
안녕하세요. 슬라이드캐스트 퀵 가이드입니다. 지금 보고 계신 이 영상은 이 가이드 파워포인트 파일 하나를 슬라이드캐스트에 넣어서 자동으로 만든 것입니다. 발표자료를 어떻게 작성해야 하는지, 영상은 어떤 과정으로 만들어지는지, 작업자가 무엇을 준비해야 하는지 차례대로 설명하겠습니다.""")
g = box(s, 0, 0, W, H, fill=BG)
box(s, 0.9, 2.15, 0.06, 2.5, fill=ACCENT)
txt(s, 1.35, 2.15, 10.5, 0.45, "SLIDECAST · QUICK GUIDE", 15, ACCENT, bold=True)
txt(s, 1.35, 2.72, 11.0, 1.3, "PPTX 한 장으로\n나레이션 영상 만들기", 46, WHITE, bold=True, space=8)
txt(s, 1.35, 4.72, 10.5, 0.5, "발표자료 작성 규칙 · 생성 프로세스 · 실행 방법 · 개발 아키텍처", 16, MUTED)
txt(s, 1.35, 6.5, 10.5, 0.4, "이 영상 자체가 이 파일로 생성되었습니다", 13, RGBColor(0x5E,0x6A,0x82), italic=True)

# ───────────────────────── 02 30초 요약
s = new("""키메시지: 30초 요약
전체 흐름은 세 단계로 아주 단순합니다. 첫째, 파워포인트의 슬라이드 노트에 발표 대본을 적습니다. 둘째, 웹 페이지에 파일을 끌어다 놓습니다. 셋째, 잠시 기다리면 음성이 입혀진 엠피포 영상과 유튜브용 자막 파일이 나옵니다. 작업자가 직접 손대야 하는 것은 사실상 첫 번째 단계, 대본 작성뿐입니다.""")
head(s, "OVERVIEW", "30초 요약 — 전체는 3단계입니다", 2)
steps = [("1", "대본 작성", "슬라이드 노트에\n발표 대본을 입력", ACCENT),
         ("2", "업로드", "웹 페이지에\n드래그 앤 드롭", AMBER),
         ("3", "결과 수령", "MP4 영상 +\nSRT 자막 다운로드", GREEN)]
for i,(n,t,b,c) in enumerate(steps):
    x = 0.9 + i*4.0
    box(s, x, 2.55, 3.6, 2.5, fill=CARD)
    cir = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.32), Inches(2.88), Inches(0.62), Inches(0.62))
    cir.fill.solid(); cir.fill.fore_color.rgb = c; cir.line.fill.background(); cir.shadow.inherit=False
    txt(s, x+0.32, 3.0, 0.62, 0.4, n, 20, BG, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x+0.32, 3.72, 3.0, 0.4, t, 20, WHITE, bold=True)
    txt(s, x+0.32, 4.18, 3.0, 0.9, b, 13, MUTED, space=3)
    if i < 2: arrow(s, x+3.72, 3.72)
box(s, 0.9, 5.35, 11.53, 0.95, fill=CARD2)
txt(s, 1.25, 5.6, 11.0, 0.5, "작업자가 실제로 손대는 것은 1단계뿐 — 나머지는 전부 자동입니다.", 16, WHITE, bold=True)

# ───────────────────────── 03 작성 가이드 ① 노트 규칙
s = new("""키메시지: 노트 규칙
가장 중요한 규칙입니다. 슬라이드 노트의 첫 줄은 나레이션에서 제외됩니다. 첫 줄은 키 메시지나 제목을 적는 자리로 쓰고, 실제로 읽히길 원하는 대본은 두 번째 줄부터 적어야 합니다. 첫 줄만 적고 끝내면 그 슬라이드는 대본이 없는 것으로 처리되어 소리 없이 지나갑니다. 이 점을 꼭 기억해 주세요.""")
head(s, "AUTHORING · 1", "가장 중요한 규칙 — 노트 첫 줄은 읽지 않습니다", 3)
box(s, 0.9, 2.4, 5.6, 3.5, fill=CARD, line=CORAL)
txt(s, 1.25, 2.68, 5.0, 0.35, "✗  이렇게 쓰면 무음 처리", 15, CORAL, bold=True)
code(s, 1.25, 3.2, 4.9, ["[슬라이드 노트]", "오늘의 핵심은 자동화입니다"], 12)
txt(s, 1.25, 4.5, 5.0, 1.2, "한 줄만 적으면 그 한 줄이\n'첫 줄'로 간주되어 통째로 제외됩니다.\n→ 나레이션 없음", 13, MUTED, space=4)

box(s, 6.83, 2.4, 5.6, 3.5, fill=CARD, line=GREEN)
txt(s, 7.18, 2.68, 5.0, 0.35, "✓  이렇게 써야 정상", 15, GREEN, bold=True)
code(s, 7.18, 3.2, 4.9, ["[슬라이드 노트]", "핵심: 자동화          ← 제외됨", "오늘의 핵심은 자동화입니다.", "세 단계로 나누어 설명하겠습니다."], 12)
txt(s, 7.18, 5.05, 5.0, 0.8, "1줄 = 키메시지(제외)\n2줄부터 = 실제 나레이션", 13, MUTED, space=4)
txt(s, 0.9, 6.15, 11.5, 0.4, "※ 2줄 이후의 여러 문단은 모두 이어서 한 번에 읽힙니다.", 13, MUTED)

# ───────────────────────── 04 작성 가이드 ② 제약
s = new("""키메시지: 작성 제약
그 밖에 지켜야 할 제약들입니다. 숨김 처리한 슬라이드는 영상에서 자동으로 빠지고 대본도 함께 제외되므로, 안 쓰는 슬라이드는 지우지 말고 숨겨두면 됩니다. 파일 크기는 백 메가바이트, 슬라이드는 백 장까지 지원합니다. 애니메이션과 화면 전환 효과는 반영되지 않고 각 슬라이드의 최종 상태만 이미지로 캡처됩니다. 영상은 가로 십육 대 구 비율로 만들어지므로 원본도 같은 비율로 작성하는 것이 좋습니다.""")
head(s, "AUTHORING · 2", "그 밖에 지켜야 할 제약", 4)
items = [("숨김 슬라이드", "숨긴 슬라이드는 영상에서 제외되고\n대본도 함께 빠집니다.\n안 쓰는 장은 삭제 대신 숨기기 권장", ACCENT),
         ("용량 · 장수 제한", "파일 100MB / 슬라이드 100장까지\n초과 시 업로드 단계에서 거부됩니다", AMBER),
         ("애니메이션 미반영", "화면 전환·애니메이션 효과는 무시되고\n각 장의 최종 상태만 캡처됩니다", CORAL),
         ("16:9 권장", "영상은 1920×1080으로 출력됩니다.\n다른 비율은 여백(레터박스)이 생깁니다", GREEN),
         ("폰트", "Mac은 Keynote로 원본 폰트를 보존합니다.\nLinux 배포 시 한글 폰트 설치 필요", ACCENT),
         ("이미지 화질", "PPT에 저해상도 이미지를 넣으면\n영상에서도 그대로 뭉개집니다", MUTED)]
for i,(t,b,c) in enumerate(items):
    x = 0.9 + (i%3)*3.93; y = 2.45 + (i//3)*1.95
    card(s, x, y, 3.63, 1.72, t, b, accent=c, tsize=15, bsize=11.5)

# ───────────────────────── 05 작업자 체크리스트
s = new("""키메시지: 작업자 체크리스트
업로드 버튼을 누르기 전에 확인할 목록입니다. 모든 슬라이드의 노트가 두 줄 이상인지, 숫자나 영문 약어를 읽는 방식이 어색하지 않은지, 문장이 너무 길지 않은지 확인해 주세요. 티티에스는 마침표와 쉼표를 기준으로 호흡을 나누기 때문에, 한 문장을 지나치게 길게 쓰면 부자연스럽게 들립니다. 두 문장에서 네 문장 정도로 끊어 쓰는 것을 권합니다.""")
head(s, "CHECKLIST", "업로드 전 작업자 확인 목록", 5)
checks = [("노트가 2줄 이상인가", "첫 줄은 제외되므로 최소 2줄 필요"),
          ("문장을 짧게 끊었는가", "한 문장 2~4개 단위 권장. TTS는 마침표로 호흡을 나눔"),
          ("숫자·약어 읽기가 자연스러운가", "'2026년' '3분기' 등은 한글로 풀어 쓰면 더 정확"),
          ("영문·특수문자를 정리했는가", "URL·기호는 소리로 읽히니 대본에서는 빼는 편이 낫다"),
          ("불필요한 슬라이드를 숨겼는가", "숨김 처리하면 영상과 대본에서 함께 제외"),
          ("목소리와 여백을 정했는가", "업로드 화면에서 성별·언어와 전환 여백을 선택")]
for i,(t,b) in enumerate(checks):
    y = 2.4 + i*0.71
    box(s, 0.9, y, 11.53, 0.6, fill=CARD if i%2==0 else CARD2)
    ck = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(y+0.14), Inches(0.32), Inches(0.32))
    ck.adjustments[0]=0.25; ck.fill.solid(); ck.fill.fore_color.rgb = ACCENT
    ck.line.fill.background(); ck.shadow.inherit=False
    txt(s, 1.15, y+0.17, 0.32, 0.3, "✓", 13, BG, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 1.68, y+0.16, 4.4, 0.35, t, 14, WHITE, bold=True)
    txt(s, 6.1, y+0.18, 6.1, 0.35, b, 12, MUTED)

# ───────────────────────── 06 개발 아키텍처
s = new("""키메시지: 개발 아키텍처
개발 구조를 보겠습니다. 사용자가 브라우저에서 파일을 올리면 패스트에이피아이 서버가 이를 받아 레디스 큐에 작업을 등록합니다. 실제 무거운 작업은 셀러리 워커가 별도 프로세스에서 처리합니다. 워커는 다섯 개의 코어 모듈을 순서대로 호출해서 최종 영상과 자막을 만듭니다. 서버와 워커를 분리했기 때문에 영상 변환이 오래 걸려도 웹 페이지는 멈추지 않습니다.""")
head(s, "ARCHITECTURE", "개발 아키텍처 — API와 워커의 분리", 6)
top = [("브라우저 UI", "static/app.js"), ("FastAPI", "src/main.py"), ("Redis 큐", "broker + backend"), ("Celery Worker", "worker/tasks.py")]
for i,(t,sub) in enumerate(top):
    x = 0.9 + i*3.05
    chip(s, x, 2.35, 2.65, 1.0, t, sub, fill=CARD2 if i<3 else CARD, edge=ACCENT if i==3 else None)
    if i < 3: arrow(s, x+2.72, 2.77, 0.26)
cores = [("pptx_parser", "노트 → 대본"), ("image_converter", "PPTX → PNG"), ("tts_engine", "대본 → MP3"),
         ("video_composer", "PNG+MP3 → MP4"), ("subtitle_maker", "타임코드 → SRT")]
box(s, 0.9, 3.75, 11.53, 1.75, fill=RGBColor(0x16,0x1B,0x27))
txt(s, 1.15, 3.92, 6.0, 0.3, "src/core/ — 워커가 순서대로 호출하는 5개 모듈", 12, ACCENT, bold=True)
for i,(t,sub) in enumerate(cores):
    x = 1.15 + i*2.2
    chip(s, x, 4.35, 2.0, 0.92, t, sub, fill=CARD2)
    if i < 4: arrow(s, x+2.04, 4.72, 0.12)
for i,(t,sub,c) in enumerate([("final_*.mp4", "1920×1080 · H.264 · AAC 48kHz", GREEN), ("subtitles_*.srt", "YouTube 업로드용 자막", AMBER)]):
    x = 0.9 + i*5.9
    box(s, x, 5.72, 5.63, 0.86, fill=CARD, line=c)
    txt(s, x+0.3, 5.9, 5.0, 0.32, t, 15, c, bold=True)
    txt(s, x+0.3, 6.24, 5.0, 0.3, sub, 11.5, MUTED)

# ───────────────────────── 07 생성 프로세스
s = new("""키메시지: 생성 프로세스
영상이 만들어지는 과정은 여섯 단계입니다. 노트를 추출하고, 슬라이드를 이미지로 바꾸고, 대본을 음성으로 합성하고, 슬라이드마다 클립을 만들고, 자막을 계산하고, 마지막으로 전체를 하나로 병합합니다. 화면의 진행률 표시가 이 여섯 단계를 그대로 따라갑니다. 진행률이 어디서 멈췄는지 보면 어느 단계에서 문제가 생겼는지 바로 알 수 있습니다.""")
head(s, "PROCESS", "생성 프로세스 — 6단계 파이프라인", 7)
steps = [("10%","노트 추출","pptx_parser","숨김 슬라이드 제외\n첫 줄 제외 후 대본화"),
         ("30%","이미지 변환","image_converter","Keynote로 PDF 추출\n→ PyMuPDF로 PNG"),
         ("30→80%","음성 합성","tts_engine","edge-tts로 슬라이드별 MP3\n길이(초)를 함께 계산"),
         ("~80%","클립 생성","video_composer","이미지+오디오를 슬라이드별\nMP4 클립으로 합성"),
         ("85%","자막 생성","subtitle_maker","오디오 길이 누적으로\nSRT 타임코드 계산"),
         ("90→100%","전체 병합","concat","클립을 무손실로 이어붙여\n최종 MP4 완성")]
for i,(pct,t,mod,b) in enumerate(steps):
    x = 0.9 + (i%3)*3.93; y = 2.4 + (i//3)*2.05
    box(s, x, y, 3.63, 1.82, fill=CARD)
    txt(s, x+0.3, y+0.24, 1.6, 0.3, pct, 12, ACCENT, bold=True)
    txt(s, x+0.3, y+0.56, 3.0, 0.35, f"{i+1}. {t}", 16, WHITE, bold=True)
    txt(s, x+0.3, y+0.94, 3.1, 0.7, b, 11.5, MUTED, space=3)
    txt(s, x+0.3, y+1.5, 3.1, 0.26, mod, 10.5, RGBColor(0x5E,0x6A,0x82))
box(s, 0.9, 6.02, 11.53, 0.46, fill=CARD2)
txt(s, 1.2, 6.13, 11.0, 0.3, "진행률이 멈춘 지점 = 문제가 생긴 단계. 30%에서 멈추면 이미지 변환, 85% 부근이면 병합 단계를 의심하세요.", 11.5, MUTED)

# ───────────────────────── 08 실행 방법
s = new("""키메시지: 실행 방법
로컬에서 실행하는 방법입니다. 먼저 에프에프엠펙과 레디스, 리브레오피스를 설치하고 파이썬 가상환경에 패키지를 깔아둡니다. 그다음 터미널 세 개를 띄워서 각각 레디스 서버, 셀러리 워커, 유비콘 웹서버를 실행합니다. 마지막으로 브라우저에서 로컬호스트 팔천 번 포트로 접속하면 업로드 화면이 나옵니다. 워커를 띄우지 않으면 업로드는 되지만 진행률이 영 퍼센트에서 멈추니 주의하세요.""")
head(s, "HOW TO RUN", "실행 방법 — 터미널 3개", 8)
txt(s, 0.9, 2.3, 6.0, 0.3, "① 최초 1회 준비", 14, ACCENT, bold=True)
code(s, 0.9, 2.68, 5.6, ["$ brew install ffmpeg redis libreoffice",
                          "$ python3 -m venv .venv",
                          "$ source .venv/bin/activate",
                          "$ pip install -r requirements.txt"])
txt(s, 0.9, 4.62, 6.0, 0.3, "② 접속", 14, ACCENT, bold=True)
code(s, 0.9, 5.0, 5.6, ["$ open http://localhost:8000"])
txt(s, 6.83, 2.3, 6.0, 0.3, "③ 매번 — 터미널 3개를 각각 실행", 14, ACCENT, bold=True)
for i,(label, cmd) in enumerate([("터미널 1 · 메시지 브로커", "$ redis-server"),
                                  ("터미널 2 · 백그라운드 워커", "$ celery -A src.worker.celery_app worker -l info"),
                                  ("터미널 3 · 웹 서버", "$ uvicorn src.main:app --reload")]):
    y = 2.68 + i*1.18
    txt(s, 6.83, y, 5.6, 0.28, label, 12, MUTED)
    code(s, 6.83, y+0.3, 5.6, [cmd], 11.5)
box(s, 6.83, 6.28, 5.6, 0.72, fill=CARD, line=AMBER)
txt(s, 7.13, 6.46, 5.1, 0.4, "⚠ 워커를 띄우지 않으면 진행률이 0%에서 멈춥니다", 12.5, AMBER, bold=True)

# ───────────────────────── 09 설정 옵션
s = new("""키메시지: 설정 옵션
업로드 화면에서 고를 수 있는 설정입니다. 목소리는 한국어 남녀와 영어 남녀를 포함해 여섯 가지를 지원합니다. 화면 전환 여백은 각 슬라이드가 바뀐 뒤 나레이션이 시작되기까지의 정적 시간으로, 영 초에서 십 초까지 지정할 수 있습니다. 기본값인 일 점 오 초는 대부분의 발표에서 자연스럽게 들립니다. 여백은 자막 타임코드에도 똑같이 반영됩니다.""")
head(s, "OPTIONS", "업로드 시 선택하는 설정", 9)
box(s, 0.9, 2.4, 6.9, 3.9, fill=CARD)
txt(s, 1.25, 2.68, 6.0, 0.35, "목소리 (edge-tts)", 17, WHITE, bold=True)
voices = [("ko_female","한국어 · 여성","SunHi"),("ko_male","한국어 · 남성","Hyunsu"),
          ("en_female","영어 · 여성","Jenny"),("en_male","영어 · 남성","Guy"),
          ("ja_male","일본어 · 남성","Keita"),("zh_male","중국어 · 남성","Yunxi")]
for i,(k,d,n) in enumerate(voices):
    y = 3.2 + i*0.5
    txt(s, 1.25, y, 1.9, 0.3, k, 12.5, ACCENT)
    txt(s, 3.3, y, 2.6, 0.3, d, 12.5, WHITE)
    txt(s, 5.9, y, 1.6, 0.3, n, 12.5, MUTED)
box(s, 8.13, 2.4, 4.3, 3.9, fill=CARD)
txt(s, 8.48, 2.68, 3.6, 0.35, "화면 전환 여백", 17, WHITE, bold=True)
txt(s, 8.48, 3.2, 3.6, 0.9, "슬라이드가 바뀐 뒤\n나레이션이 시작되기까지의\n정적 시간", 12.5, MUTED, space=3)
box(s, 8.48, 4.25, 3.6, 0.85, fill=CARD2)
txt(s, 8.48, 4.45, 3.6, 0.45, "1.5초", 24, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 8.48, 5.25, 3.6, 0.8, "허용 범위 0 ~ 10초\n자막 타임코드에도\n동일하게 반영됩니다", 12, MUTED, space=3)
txt(s, 0.9, 6.38, 11.5, 0.35, "설정을 바꾸면 같은 PPTX로도 다른 결과물을 만들 수 있습니다. 언어별 대본은 노트를 바꿔 다시 업로드하세요.", 12.5, MUTED)

# ───────────────────────── 10 산출물
s = new("""키메시지: 산출물
작업이 끝나면 두 개의 파일이 나옵니다. 하나는 풀에이치디 해상도의 엠피포 영상이고, 다른 하나는 유튜브에 그대로 올릴 수 있는 에스알티 자막입니다. 자막의 타임코드는 각 슬라이드의 실제 음성 길이를 누적해서 계산하기 때문에 영상과 정확히 맞습니다. 영상은 다운로드 버튼으로 받고, 자막 파일은 작업 폴더에서 함께 확인할 수 있습니다.""")
head(s, "OUTPUT", "산출물 — MP4와 SRT", 10)
box(s, 0.9, 2.45, 5.75, 3.5, fill=CARD, line=GREEN)
txt(s, 1.25, 2.75, 5.0, 0.4, "final_{job_id}.mp4", 19, GREEN, bold=True)
for i,(k,v) in enumerate([("해상도","1920 × 1080 (Full HD)"),("프레임","30 fps"),
                          ("비디오","H.264 / CRF 18"),("오디오","AAC 192kbps · 48kHz 스테레오"),
                          ("길이","각 슬라이드 = 음성 길이 + 전환 여백")]):
    y = 3.4 + i*0.48
    txt(s, 1.25, y, 1.5, 0.3, k, 12, MUTED)
    txt(s, 2.85, y, 3.6, 0.3, v, 12.5, WHITE)
box(s, 7.0, 2.45, 5.43, 3.5, fill=CARD, line=AMBER)
txt(s, 7.35, 2.75, 5.0, 0.4, "subtitles_{job_id}.srt", 19, AMBER, bold=True)
txt(s, 7.35, 3.4, 4.8, 0.6, "슬라이드별 음성 길이를 누적해\n타임코드를 계산 — 영상과 정확히 일치", 12.5, MUTED, space=3)
code(s, 7.35, 4.2, 4.75, ["1", "00:00:01,500 --> 00:00:16,548", "안녕하세요. 슬라이드캐스트..."], 11)
txt(s, 7.35, 5.5, 4.8, 0.4, "대본이 없는 슬라이드는 자막에서도 제외", 11.5, MUTED)
box(s, 0.9, 6.14, 11.53, 0.5, fill=CARD2)
txt(s, 1.25, 6.25, 11.0, 0.4, "영상은 완료 화면의 다운로드 버튼으로, 자막은 작업 폴더에서 함께 받을 수 있습니다.", 12, WHITE)

# ───────────────────────── 11 트러블슈팅
s = new("""키메시지: 트러블슈팅
자주 나오는 증상과 원인입니다. 특정 슬라이드가 소리 없이 지나간다면 노트를 한 줄만 적었을 가능성이 큽니다. 진행률이 영 퍼센트에서 멈춘다면 셀러리 워커가 실행되지 않은 것입니다. 한글이 깨져 보인다면 리눅스 서버에 한글 폰트가 없는 경우입니다. 대본과 화면이 어긋난다면 예전 버전을 쓰고 있는 것이니 최신 코드로 갱신해 주세요.""")
head(s, "TROUBLESHOOTING", "이런 증상이 나오면", 11)
rows = [("특정 슬라이드가 무음으로 지나감", "노트를 1줄만 작성 — 첫 줄은 제외됩니다", "노트를 2줄 이상으로 수정", CORAL),
        ("진행률이 0%에서 멈춤", "Celery 워커가 실행되지 않음", "터미널 2에서 워커 실행 확인", AMBER),
        ("업로드가 즉시 거부됨", "100MB 초과 / .pptx가 아닌 파일 / 잘못된 설정값", "화면에 표시된 사유 확인", AMBER),
        ("한글이 네모(□)로 깨짐", "Linux 서버에 한글 폰트 미설치", "서버에 나눔고딕 등 설치 (Mac은 무관)", ACCENT),
        ("대본과 화면이 어긋남", "숨김 슬라이드 정렬 버그 (구버전)", "최신 코드로 갱신 — Phase 5.1에서 수정됨", GREEN)]
txt(s, 1.2, 2.3, 4.0, 0.3, "증상", 12, RGBColor(0x5E,0x6A,0x82), bold=True)
txt(s, 5.3, 2.3, 3.7, 0.3, "원인", 12, RGBColor(0x5E,0x6A,0x82), bold=True)
txt(s, 9.2, 2.3, 3.2, 0.3, "조치", 12, RGBColor(0x5E,0x6A,0x82), bold=True)
for i,(sym,cause,fix,c) in enumerate(rows):
    y = 2.68 + i*0.82
    box(s, 0.9, y, 11.53, 0.72, fill=CARD if i%2==0 else CARD2)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(y+0.14), Inches(0.045), Inches(0.44))
    bar.fill.solid(); bar.fill.fore_color.rgb = c; bar.line.fill.background(); bar.shadow.inherit=False
    txt(s, 1.2, y+0.24, 4.0, 0.32, sym, 13, WHITE, bold=True)
    txt(s, 5.3, y+0.26, 3.8, 0.32, cause, 11.5, MUTED)
    txt(s, 9.2, y+0.26, 3.1, 0.32, fix, 11.5, MUTED)

# ───────────────────────── 12 최근 안정화
s = new("""키메시지: 최근 안정화
최근에 두 가지 치명적인 버그를 고쳤습니다. 첫째, 대본이 없는 슬라이드가 있으면 그 이후 나레이션이 통째로 앞당겨지는 문제가 있었습니다. 둘째, 숨김 슬라이드가 있으면 대본과 화면이 어긋나는 문제가 있었습니다. 두 버그 모두 에러 없이 조용히 잘못된 결과물을 만들어냈기 때문에 특히 위험했습니다. 지금은 모두 수정되어 이 영상처럼 정상적으로 동작합니다.""")
head(s, "STABILITY", "최근 수정된 치명 버그 (Phase 5.1)", 12)
for i,(t,before,after) in enumerate([
    ("무음 슬라이드 병합 파손",
     "대본 없는 슬라이드의 오디오 트랙이 누락되어\n이후 나레이션이 통째로 앞당겨짐 (영상 64.3s / 오디오 62.9s)",
     "무음 트랙을 삽입하고 전 클립의 오디오 규격을 통일\n→ 영상 64.2s / 오디오 64.3s 일치"),
    ("슬라이드-대본 정렬 어긋남",
     "숨김 슬라이드가 대본에는 포함되고 이미지에는 빠져\n그 이후 전부 다른 슬라이드의 나레이션이 붙음",
     "파서에서 숨김 슬라이드를 제외하고\n개수 불일치 시 조용히 진행하지 않고 즉시 에러")]):
    y = 2.32 + i*2.08
    box(s, 0.9, y, 11.53, 1.85, fill=CARD)
    txt(s, 1.25, y+0.24, 6.0, 0.35, f"🔴 {t}", 17, WHITE, bold=True)
    txt(s, 1.25, y+0.72, 5.3, 0.9, "BEFORE", 11, CORAL, bold=True)
    txt(s, 1.25, y+0.98, 5.3, 0.8, before, 11.5, MUTED, space=2)
    txt(s, 6.95, y+0.72, 5.3, 0.9, "AFTER", 11, GREEN, bold=True)
    txt(s, 6.95, y+0.98, 5.3, 0.8, after, 11.5, MUTED, space=2)
txt(s, 0.9, 6.36, 11.5, 0.35, "두 버그 모두 에러 없이 잘못된 결과물을 만들어냈기 때문에 특히 위험했습니다.", 12.5, AMBER)

# ───────────────────────── 13 남은 로드맵
s = new("""키메시지: 남은 로드맵
앞으로 남은 작업입니다. 티티에스를 병렬로 처리해서 생성 속도를 높이는 작업과, 임시 파일을 자동으로 정리하는 작업이 다음 순서입니다. 그 뒤로는 회원과 결제, 클라우드 배포, 유튜브 자동 업로드가 이어집니다. 지금 단계에서도 로컬에서 영상을 만드는 데는 아무 문제가 없습니다.""")
head(s, "ROADMAP", "다음에 할 일", 13)
todo = [("다음","TTS 병렬화","슬라이드마다 순차 호출하는 음성 생성을\n동시 요청으로 전환 — 장수에 비례해 단축", GREEN),
        ("다음","임시 파일 자동 정리","완료 후 PNG·MP3·개별 클립을 삭제하고\n최종 영상과 자막만 보존", GREEN),
        ("이후","WebSocket 실연결","현재 진행률은 폴링으로만 동작.\nRedis pub/sub 연결 또는 코드 정리", AMBER),
        ("이후","회원 · 결제","이메일 로그인, 무료/유료 플랜,\n결제 연동", MUTED),
        ("이후","클라우드 배포","API 서버와 워커 서버 분리,\nS3 연동, 한글 폰트 자동 설치", MUTED),
        ("이후","YouTube 연동","다국어 자막, 챕터, 썸네일,\nAPI 자동 업로드", MUTED)]
for i,(tag,t,b,c) in enumerate(todo):
    x = 0.9 + (i%3)*3.93; y = 2.45 + (i//3)*1.95
    card(s, x, y, 3.63, 1.72, t, b, accent=c, tsize=15, bsize=11.5)
    txt(s, x+3.63-0.75, y+0.32, 0.5, 0.25, tag, 10, c, bold=True, align=PP_ALIGN.RIGHT)

# ───────────────────────── 14 마무리
s = new("""키메시지: 마무리
정리하겠습니다. 노트 두 번째 줄부터가 대본이고, 나머지는 전부 자동입니다. 지금 보신 이 영상은 이 가이드 파일 하나만으로 만들어졌습니다. 여러분의 발표자료도 같은 방식으로 만들어 보세요. 감사합니다.""")
box(s, 0.9, 1.9, 0.06, 3.6, fill=ACCENT)
txt(s, 1.35, 1.9, 11.0, 0.45, "SUMMARY", 15, ACCENT, bold=True)
txt(s, 1.35, 2.45, 11.0, 1.2, "노트 2번째 줄부터가 대본,\n나머지는 전부 자동입니다.", 40, WHITE, bold=True, space=8)
for i,(k,v) in enumerate([("작성","슬라이드 노트에 대본 (첫 줄 제외)"),
                          ("실행","redis-server · celery worker · uvicorn"),
                          ("결과","1920×1080 MP4 + YouTube용 SRT")]):
    y = 4.35 + i*0.55
    txt(s, 1.35, y, 1.2, 0.3, k, 13, ACCENT, bold=True)
    txt(s, 2.6, y, 9.5, 0.3, v, 14, MUTED)
txt(s, 1.35, 6.35, 11.0, 0.4, "이 영상은 이 가이드 PPTX 파일 하나로 생성되었습니다.", 14, RGBColor(0x5E,0x6A,0x82), italic=True)

prs.save(sys.argv[1])
print(f"생성 완료: {sys.argv[1]}  ({len(prs.slides.__iter__.__self__._sldIdLst)}장)")
