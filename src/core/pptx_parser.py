from pptx import Presentation
from typing import List, Dict, Any


def _is_hidden(slide) -> bool:
    """
    PPTX에서 '슬라이드 숨기기'가 적용된 슬라이드인지 판별합니다.
    OOXML상 <p:sld show="0"> 속성으로 표현된다.
    """
    show = slide.element.get("show")
    return show is not None and show in ("0", "false")


def extract_notes(pptx_path: str, include_hidden: bool = False) -> List[Dict[str, Any]]:
    """
    PPTX 파일에서 슬라이드 번호와 노트(발표 스크립트)를 추출합니다.

    기본적으로 숨김 슬라이드는 제외한다. Keynote/LibreOffice의 PDF 내보내기가
    숨김 슬라이드를 건너뛰기 때문에, 여기서 포함해버리면 대본과 슬라이드 이미지가
    통째로 밀려 잘못된 나레이션이 붙는다.

    반환 항목:
        slide_no    - 렌더링되는 슬라이드 순번(1-based). slide_001.png 와 1:1 대응
        pptx_index  - 원본 PPTX에서의 순번(1-based). 숨김 슬라이드 포함 기준
        hidden      - 숨김 슬라이드 여부
        notes       - 발표 대본 (노트의 첫 줄은 제목/키메시지로 보고 제외)
    """
    prs = Presentation(pptx_path)
    slides = []

    for src_idx, slide in enumerate(prs.slides):
        hidden = _is_hidden(slide)
        if hidden and not include_hidden:
            continue

        notes = ""
        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if text_frame:
                raw_notes = text_frame.text.strip()
                if raw_notes:
                    # 첫 줄(제목/키메시지) 제외하고 나머지 전체를 스크립트로 사용
                    lines = raw_notes.split('\n', 1)
                    notes = lines[1].strip() if len(lines) > 1 else ""

        slides.append({
            "slide_no": len(slides) + 1,
            "pptx_index": src_idx + 1,
            "hidden": hidden,
            "notes": notes
        })

    return slides
